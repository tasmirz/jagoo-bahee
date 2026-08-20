from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("Islands-of-Reach-Presentation.pptx")

INK = "1F2933"
BLUE = "2C6E9B"
TEAL = "3F8F7A"
AMBER = "C6873B"
ROSE = "A8465F"
MIST = "DDE3E8"
PAPER = "F7F6F2"
WHITE = "FFFFFF"
MUTED = "5B6770"


def rgb(value):
    return RGBColor.from_string(value)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_background(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.02, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.font.color.rgb = rgb(color)
    return shape


def add_rich_text(slide, runs, x, y, w, h, size=20, color=INK,
                  valign=MSO_ANCHOR.TOP, margin=0.08):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    for item in runs:
        run = p.add_run()
        run.text = item[0]
        run.font.name = "Aptos"
        run.font.size = Pt(item[1] if len(item) > 1 else size)
        run.font.bold = item[2] if len(item) > 2 else False
        run.font.color.rgb = rgb(item[3] if len(item) > 3 else color)
    return shape


def add_box(slide, text, x, y, w, h, fill=WHITE, line=MIST, size=17,
            bold=False, color=INK, radius=True, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE, margin=0.1):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def add_rule(slide, x1, y1, x2, y2, color=MIST, width=1.2, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = 4
    return line


def add_chevron(slide, x, y, w=0.22, h=0.32, color=MUTED):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_slide(title, kicker=None, section_color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    if kicker:
        add_text(slide, kicker.upper(), 0.65, 0.34, 4.5, 0.28, 10,
                 section_color, True)
    add_text(slide, title, 0.65, 0.67 if kicker else 0.48, 12.0, 0.58,
             28, INK, True)
    add_rule(slide, 0.65, 1.27, 12.68, 1.27, MIST, 1)
    return slide


def add_footer(slide, number, source=None):
    if source:
        add_text(slide, source, 0.65, 7.13, 11.3, 0.18, 8.5, MUTED)
    add_text(slide, str(number), 12.28, 7.08, 0.4, 0.2, 9, MUTED,
             align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=20, color=INK, gap=7):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
    return shape


def label_box(slide, heading, body, x, y, w, h, fill=WHITE, line=MIST,
              heading_color=INK, body_size=15):
    box = add_box(slide, "", x, y, w, h, fill, line, radius=True)
    add_text(slide, heading, x + 0.18, y + 0.13, w - 0.36, 0.33, 16,
             heading_color, True)
    add_text(slide, body, x + 0.18, y + 0.52, w - 0.36, h - 0.64,
             body_size, INK)
    return box


# 1. Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_box(slide, "", 0, 0, 0.22, 7.5, BLUE, BLUE, radius=False)
add_text(slide, "Islands of Reach", 0.85, 1.25, 9.8, 0.75, 34, INK, True)
add_text(slide, "A censorship-resistant federated forum for national internet blackouts",
         0.88, 2.12, 9.2, 1.0, 24, BLUE)
add_text(slide, "Signed public discussion across the IP paths that still exist",
         0.88, 3.34, 7.8, 0.4, 17, MUTED)
add_rule(slide, 0.88, 4.05, 6.4, 4.05, TEAL, 3)
add_text(slide, "NSysS 2026", 0.88, 6.45, 3.0, 0.3, 13, MUTED, True)
add_text(slide, "Cox’s Bazar, Bangladesh", 0.88, 6.78, 4.0, 0.26, 12, MUTED)
add_footer(slide, 1)


# 2. Problem
slide = add_slide("When the outside internet disappears", "Problem statement", ROSE)
add_text(slide,
         "A public forum normally depends on distant platforms, identity providers and globally reachable servers.",
         0.72, 1.55, 7.05, 0.8, 23, INK, True)
steps = [
    ("1", "Traffic is throttled", "Messages become slow and unreliable"),
    ("2", "Services are blocked", "Circumvention still needs a path outside"),
    ("3", "International transit is cut", "Domestic networks may survive as separate islands"),
]
for i, (n, head, body) in enumerate(steps):
    y = 3.05 + i * 1.12
    add_box(slide, n, 0.78, y, 0.52, 0.52, ROSE, ROSE, 18, True, WHITE)
    add_text(slide, head, 1.48, y - 0.01, 3.0, 0.32, 18, INK, True)
    add_text(slide, body, 1.48, y + 0.36, 5.5, 0.34, 14, MUTED)
    if i < 2:
        add_rule(slide, 1.04, y + 0.55, 1.04, y + 1.12, ROSE, 1.5)
label_box(slide, "The design question",
          "Can independent domestic servers keep a readable, verifiable forum working across the remaining ISP links?",
          8.25, 1.65, 4.25, 3.25, "F3E8EB", ROSE, ROSE, 20)
add_text(slide, "Scope: surviving IP connectivity. This is not mesh or off-grid networking.",
         8.38, 5.28, 4.0, 0.7, 17, INK, True)
add_footer(slide, 2, "Sources: OONI (2025); Bischof et al. (2025); Telex (USENIX Security 2011).")


# 3. Abstract
slide = add_slide("The solution in one minute", "Simplified abstract", TEAL)
summary = [
    ("Independent servers", "Communities can run their own nodes and choose their federation peers."),
    ("Signed objects", "A post is valid because its author signed it, not because one server approved it."),
    ("Portable evidence", "A client stores the node receipt and can copy it to independent audit-log servers."),
    ("Bridged ISP islands", "A trusted multi-homed node relays verified objects across surviving domestic paths."),
]
for i, (head, body) in enumerate(summary):
    x = 0.72 + (i % 2) * 6.08
    y = 1.55 + (i // 2) * 2.05
    color = [BLUE, TEAL, AMBER, ROSE][i]
    label_box(slide, head, body, x, y, 5.62, 1.62, WHITE, color, color, 17)
add_box(slide,
        "If no path remains between two components, each side keeps working locally and queues data. Crossing resumes only when a bridge can reach both sides.",
        1.32, 5.83, 10.7, 0.82, "EEF3F6", BLUE, 17, True, INK)
add_footer(slide, 3)


# 4. Ladder
slide = add_slide("Resilience is a ladder, not a single claim", "Reachability model", BLUE)
levels = [
    ("L0", "Ordinary reach", "Global paths are available"),
    ("L1", "National partition", "International transit is withdrawn"),
    ("L2", "ISP island", "Domestic peering is also cut"),
    ("L3", "Bridged islands", "A multi-homed operator joins two islands"),
]
for i, (level, name, detail) in enumerate(levels):
    y = 1.52 + i * 1.28
    color = [MIST, "D9E9E4", "DCE9F1", "F3E5D5"][i]
    border = [MUTED, TEAL, BLUE, AMBER][i]
    add_box(slide, level, 0.85, y, 0.78, 0.78, border, border, 19, True, WHITE)
    add_box(slide, "", 1.8, y, 7.1, 0.78, color, border, radius=True)
    add_text(slide, name, 2.05, y + 0.12, 2.6, 0.27, 18, INK, True)
    add_text(slide, detail, 4.53, y + 0.13, 4.05, 0.3, 15, MUTED)
    if i < 3:
        add_rule(slide, 1.24, y + 0.8, 1.24, y + 1.24, border, 1.5)
label_box(slide, "What is measured",
          "The prototype exercises L0 through L3. It does not promise delivery when every physical link is gone.",
          9.48, 2.0, 3.1, 2.55, "F5EEE4", AMBER, AMBER, 17)
add_footer(slide, 4)


# 5. Related work
slide = add_slide("Related systems solve parts of the problem", "Literature review", AMBER)
rows = [
    ("Publius, Freenet, Tangler", "Host takedown and deletion", "Readers still need a route to replicas"),
    ("Tor, Snowflake, Telex", "Blocking and circumvention", "They require some path beyond the censor"),
    ("Netnews, ActivityPub, Matrix", "Independent servers and public discussion", "No scope-aware bridge with this evidence path"),
    ("Nostr, SSB, AT Protocol", "Signed or content-addressed objects", "Different relay, identity or partition assumptions"),
]
headers = [("System family", 0.72, 3.3), ("What it gives", 4.18, 3.3), ("Gap for this setting", 7.68, 4.88)]
for text_value, x, w in headers:
    add_box(slide, text_value, x, 1.48, w, 0.55, INK, INK, 15, True, WHITE, radius=False)
for i, row in enumerate(rows):
    y = 2.08 + i * 0.93
    fill = WHITE if i % 2 == 0 else "EEF1F3"
    add_box(slide, row[0], 0.72, y, 3.3, 0.8, fill, MIST, 15, True, INK, radius=False, align=PP_ALIGN.LEFT)
    add_box(slide, row[1], 4.18, y, 3.3, 0.8, fill, MIST, 14, False, INK, radius=False, align=PP_ALIGN.LEFT)
    add_box(slide, row[2], 7.68, y, 4.88, 0.8, fill, MIST, 14, False, INK, radius=False, align=PP_ALIGN.LEFT)
add_text(slide,
         "The contribution is the combination: server federation, no-inbound participation, gossiped evidence and operation across L0–L3.",
         0.85, 6.08, 11.8, 0.55, 17, BLUE, True, align=PP_ALIGN.CENTER)
add_footer(slide, 5, "Selected sources: Waldman et al. (2000, 2001); Clarke et al. (2001); RFC 5537; RFC 9162; W3C ActivityPub; Tarr et al. (2019).")


# 6. Signed content
slide = add_slide("Step 1: the author signs the post", "Method", BLUE)
boxes = [
    ("Write", "Post, vote, moderation action or revocation", BLUE),
    ("Encode", "One canonical byte representation", TEAL),
    ("Sign", "Ed25519 signature covers those bytes", AMBER),
    ("Identify", "SHA-256 content ID comes from the same bytes", ROSE),
]
for i, (head, body, color) in enumerate(boxes):
    x = 0.62 + i * 3.13
    label_box(slide, head, body, x, 2.15, 2.72, 1.72, WHITE, color, color, 15)
    if i < 3:
        add_chevron(slide, x + 2.82, 2.86, 0.2, 0.3, MUTED)
add_box(slide, "Every feature enters through the same write endpoint.",
        2.0, 4.45, 4.2, 0.72, "E8F1F6", BLUE, 18, True, INK)
add_box(slide, "A relay cannot alter or re-encode the signed bytes.",
        7.12, 4.45, 4.2, 0.72, "E7F0ED", TEAL, 18, True, INK)
add_text(slide, "No email address, phone number or password is required for a pseudonym.",
         1.2, 5.75, 10.9, 0.5, 21, INK, True, align=PP_ALIGN.CENTER)
add_footer(slide, 6, "Methods: RFC 8032 Ed25519; deterministic protobuf profile; cross-language conformance vectors.")


# 7. Pipeline
slide = add_slide("Step 2: every write passes the same 19 checks", "Method", TEAL)
groups = [
    ("1–12", "Reject before storage", "size, canonical parse, domain, clock, signature, certificate, dedupe, replay", TEAL, 5.0),
    ("13–15", "Admission and meaning", "anti-abuse, authorisation, body validation", MUTED, 2.5),
    ("16–17", "One transaction", "apply projection + append Merkle witness", ROSE, 2.4),
    ("18–19", "Acknowledge and relay", "signed receipt + federation fanout", AMBER, 2.4),
]
x = 0.55
for idx, (nums, head, body, color, width) in enumerate(groups):
    add_box(slide, nums, x, 1.68, width, 0.52, color, color, 17, True, WHITE, radius=False)
    add_box(slide, "", x, 2.2, width, 2.0, WHITE, color, radius=False)
    add_text(slide, head, x + 0.18, 2.43, width - 0.36, 0.35, 17, color, True)
    add_text(slide, body, x + 0.18, 2.96, width - 0.36, 0.9, 14.5, INK)
    x += width + 0.1
add_box(slide,
        "Federated objects repeat every applicable check. Only the origin-specific payment at step 13 is not charged twice.",
        0.85, 4.72, 11.65, 0.78, "EEF3F6", BLUE, 17, True, INK)
add_bullets(slide, [
    "Steps 1–12 cannot turn invalid input into database write load.",
    "Steps 16–17 prevent a visible post from existing without a matching log entry.",
], 1.25, 5.82, 10.7, 0.86, 16)
add_footer(slide, 7)


# 8. Federation
slide = add_slide("Step 3: independent servers exchange verified objects", "Method", BLUE)
for x, label, color in [(0.85, "Server S1", TEAL), (5.25, "Server S2", BLUE), (9.65, "Server S3", AMBER)]:
    add_box(slide, label, x, 2.05, 2.75, 1.15, WHITE, color, 20, True, color)
add_rule(slide, 3.6, 2.62, 5.25, 2.62, TEAL, 2)
add_rule(slide, 8.0, 2.62, 9.65, 2.62, BLUE, 2)
add_text(slide, "stream + backfill", 3.72, 2.2, 1.42, 0.25, 11, MUTED, True, align=PP_ALIGN.CENTER)
add_text(slide, "stream + backfill", 8.12, 2.2, 1.42, 0.25, 11, MUTED, True, align=PP_ALIGN.CENTER)
items = [
    ("Verify again", "Trust changes quota, never validity."),
    ("Keep exact bytes", "A peer cannot repair malformed encoding."),
    ("Work behind CGNAT", "Outbound drains and requested backfill need no inbound address."),
    ("Recover after gaps", "Durable queues and resumable backfill continue after a path returns."),
]
for i, (head, body) in enumerate(items):
    x = 0.72 + (i % 2) * 6.08
    y = 4.05 + (i // 2) * 1.2
    label_box(slide, head, body, x, y, 5.62, 1.15, WHITE, MIST, BLUE, 14)
add_footer(slide, 8)


# 9. ALS
slide = add_slide("Step 4: the audit-log server preserves the receipt", "ALS workflow", TEAL)
top = [
    ("1. Client", "queue signed request", BLUE),
    ("2. Node", "validate and commit", TEAL),
    ("3. Receipt", "tree head + proof", AMBER),
    ("4. Certificate", "exact request + receipt\nsave locally first", BLUE),
]
for i, (head, body, color) in enumerate(top):
    x = 0.62 + i * 3.12
    label_box(slide, head, body, x, 1.55, 2.72, 1.22, WHITE, color, color, 14)
    if i < 3:
        add_chevron(slide, x + 2.82, 2.02, 0.2, 0.3, MUTED)
add_rule(slide, 10.1, 2.77, 10.1, 3.36, TEAL, 2)
add_rule(slide, 10.1, 3.36, 4.05, 3.36, TEAL, 2)
add_rule(slide, 4.05, 3.36, 4.05, 3.75, TEAL, 2)
add_rule(slide, 10.1, 3.36, 8.33, 3.36, TEAL, 2)
add_rule(slide, 8.33, 3.36, 8.33, 3.75, TEAL, 2)
label_box(slide, "ALS A", "Recompute IDs, verify signatures and inclusion proof, append hash chain.",
          2.45, 3.75, 3.2, 1.45, "E7F0ED", TEAL, TEAL, 14)
label_box(slide, "ALS B", "Run the same checks and keep an independent copy.",
          6.73, 3.75, 3.2, 1.45, "E7F0ED", TEAL, TEAL, 14)
label_box(slide, "Issuing node", "/status reports online, hidden, deleted or wrong-server status.",
          10.18, 3.75, 2.45, 1.45, "F5EEE4", AMBER, AMBER, 13)
add_box(slide,
        "The ALS proves prior acknowledgement. It does not moderate, report live state, guarantee storage or recover a body after every copy is gone.",
        0.85, 5.72, 11.65, 0.72, "F3E8EB", ROSE, 16, True, INK)
add_text(slide, "Direct-message certificates stay local to avoid creating a timestamped social graph.",
         1.1, 6.56, 11.15, 0.3, 13, MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 9)


# 10. Moderation and identity
slide = add_slide("Step 5: visible moderation, local identity", "Method", ROSE)
label_box(slide, "Moderation", "A moderator signs an additive action. Clients can show the policy and its reason.",
          0.75, 1.55, 3.75, 1.42, "F3E8EB", ROSE, ROSE, 17)
label_box(slide, "Deletion", "A signed tombstone records deletion. Receipts still prove prior acknowledgement.",
          4.78, 1.55, 3.75, 1.42, "F5EEE4", AMBER, AMBER, 17)
label_box(slide, "Choice", "A user may move to another server with a different policy and still read federated content.",
          8.81, 1.55, 3.75, 1.42, "E8F1F6", BLUE, BLUE, 17)
add_text(slide, "One device seed", 0.95, 3.67, 2.1, 0.4, 18, INK, True)
add_box(slide, "Community A key", 3.35, 3.28, 2.35, 0.62, WHITE, TEAL, 15, True, TEAL)
add_box(slide, "Community B key", 3.35, 4.1, 2.35, 0.62, WHITE, BLUE, 15, True, BLUE)
add_box(slide, "Community C key", 3.35, 4.92, 2.35, 0.62, WHITE, AMBER, 15, True, AMBER)
for y in [3.59, 4.41, 5.23]:
    add_rule(slide, 2.78, 3.85, 3.35, y, MUTED, 1.3)
label_box(slide, "Anti-abuse without civil identity",
          "Proof of work, credits, blind credentials and epoch nullifiers price posting. They limit rate, not the number of identities.",
          6.55, 3.48, 5.45, 1.9, "EEF1F3", MIST, INK, 18)
add_text(slide, "Honest limit: network observers can still correlate timing and traffic volume.",
         6.65, 5.73, 5.2, 0.6, 17, ROSE, True)
add_footer(slide, 10, "Sources: Dwork and Naor (1993); Chaum (1983); RFC 9576; UN Special Rapporteur on anonymity (2015).")


# 11. Bridge
slide = add_slide("Bridge two ISP islands", "Method step 6", AMBER)
add_box(slide, "ISP island A\nAS 64501", 0.78, 2.0, 2.35, 1.65, "E7F0ED", TEAL, 20, True, INK)
add_box(slide, "Bridge\neth0 → island A\neth1 → island B", 5.02, 1.75, 3.3, 2.15,
        "F5EEE4", AMBER, 19, True, INK)
add_box(slide, "ISP island B\nAS 64502", 10.18, 2.0, 2.35, 1.65, "E8F1F6", BLUE, 20, True, INK)
add_rule(slide, 3.13, 2.82, 5.02, 2.82, TEAL, 2.5)
add_rule(slide, 8.32, 2.82, 10.18, 2.82, BLUE, 2.5)
add_text(slide, "direct exchange cut", 5.11, 1.29, 3.1, 0.24, 14, ROSE, True, align=PP_ALIGN.CENTER)
add_rule(slide, 2.0, 1.57, 11.3, 1.57, MUTED, 1.2, True)
add_text(slide, "×", 6.42, 1.29, 0.45, 0.46, 27, ROSE, True, align=PP_ALIGN.CENTER)
requirements = [
    "Trusted on both sides",
    "Outbound sockets bound to the correct source interface",
    "Quota counted across the uplink pair",
    "Half the grant reserved so bulk forum traffic cannot starve urgent classes",
]
add_bullets(slide, requirements, 1.0, 4.47, 11.4, 1.72, 17, INK, 5)
add_box(slide, "A bridge is a relay policy after verification, not a twentieth validation step.",
        2.0, 6.3, 9.3, 0.55, "EEF3F6", BLUE, 15, True, INK)
add_footer(slide, 11)


# 12. Crossing and failover
slide = add_slide("After the cut, dependencies cross in order", "Measured crossing", TEAL)
add_text(slide, "Island A", 0.82, 1.55, 1.4, 0.3, 16, TEAL, True)
add_text(slide, "Verifying bridge", 5.45, 1.55, 2.2, 0.3, 16, AMBER, True, align=PP_ALIGN.CENTER)
add_text(slide, "Island B", 11.05, 1.55, 1.4, 0.3, 16, BLUE, True, align=PP_ALIGN.RIGHT)
events = [
    ("Author certificate", "0.5 s", 0.5),
    ("Community record", "1.2 s", 1.2),
    ("Post", "2.1 s", 2.1),
]
for i, (name, duration, val) in enumerate(events):
    y = 2.18 + i * 0.95
    add_text(slide, name, 0.82, y + 0.12, 2.2, 0.28, 15, INK, True)
    add_box(slide, "", 3.02, y, 7.7, 0.52, "EEF1F3", MIST, radius=False)
    add_box(slide, "", 3.02, y, 7.7 * val / 2.1, 0.52, TEAL if i < 2 else BLUE,
            TEAL if i < 2 else BLUE, radius=False)
    add_text(slide, duration, 10.92, y + 0.1, 1.05, 0.3, 16, TEAL, True)
add_box(slide, "First deployed run", 1.05, 5.02, 2.5, 0.5, ROSE, ROSE, 15, True, WHITE)
add_text(slide, "407.6 s", 1.05, 5.62, 2.5, 0.5, 26, ROSE, True, align=PP_ALIGN.CENTER)
add_box(slide, "Cause", 4.15, 5.02, 2.5, 0.5, MUTED, MUTED, 15, True, WHITE)
add_text(slide, "one blackholed peer blocked a serial drain", 3.82, 5.62, 3.15, 0.72, 16, INK, align=PP_ALIGN.CENTER)
add_box(slide, "After the fix", 7.75, 5.02, 2.5, 0.5, TEAL, TEAL, 15, True, WHITE)
add_text(slide, "2.1 s", 7.75, 5.62, 2.5, 0.5, 26, TEAL, True, align=PP_ALIGN.CENTER)
add_text(slide, "Concurrent per-peer drains with deadlines", 10.43, 5.62, 2.15, 0.72, 15, INK, align=PP_ALIGN.CENTER)
add_footer(slide, 12)


# 13. Threats
slide = add_slide("What the design resists, and what it does not", "Threat model", ROSE)
add_box(slide, "Resisted or bounded", 0.72, 1.48, 5.85, 0.53, TEAL, TEAL, 17, True, WHITE, radius=False)
add_box(slide, "Not solved by this design", 6.76, 1.48, 5.85, 0.53, ROSE, ROSE, 17, True, WHITE, radius=False)
left = [
    "Post modification: author signature and content ID",
    "Deletion after acknowledgement: client and ALS retain receipts",
    "Opaque moderation: signed actions and visible local policy",
    "Invalid-input write amplification: no-write validation prefix",
    "One failed ISP path: ranked paths, bridge and durable retry",
    "Peer equivocation: signed, gossiped tree heads",
]
right = [
    "No traffic-analysis or stylometry resistance",
    "No safety after endpoint or signing-key compromise",
    "No guarantee if every link or independent copy disappears",
    "No identity-based Sybil prevention; posting cost only bounds rate",
    "No fair or universal moderation policy",
    "No guarantee against a compromised majority of peers or auditors",
]
add_bullets(slide, left, 0.82, 2.2, 5.55, 3.85, 16, INK, 8)
add_bullets(slide, right, 6.86, 2.2, 5.55, 3.85, 16, INK, 8)
add_box(slide, "Condition for crossing: at least one authorised bridge must still reach both components.",
        1.45, 6.28, 10.45, 0.55, "F5EEE4", AMBER, 16, True, INK)
add_footer(slide, 13)


# 14. Results
slide = add_slide("Prototype results", "Evaluation", BLUE)
metrics = [
    ("16 / 16", "canonical vectors agree", "TypeScript, Rust, Python", TEAL),
    ("0 loss", "200 posts reached hop 7", "median 4.125 s", BLUE),
    ("2.1 s", "cut-path post rendered", "after dependencies crossed", AMBER),
    ("730 req/s", "invalid flood rejected", "zero database writes", ROSE),
]
for i, (value, label, note, color) in enumerate(metrics):
    x = 0.65 + i * 3.13
    add_box(slide, "", x, 1.5, 2.85, 1.65, WHITE, color, radius=True)
    add_text(slide, value, x + 0.12, 1.72, 2.61, 0.43, 25, color, True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, 2.24, 2.61, 0.3, 14, INK, True, align=PP_ALIGN.CENTER)
    add_text(slide, note, x + 0.12, 2.64, 2.61, 0.25, 11.5, MUTED, align=PP_ALIGN.CENTER)

chart_data = ChartData()
chart_data.categories = ["1", "2", "3", "4", "5", "6", "7"]
chart_data.add_series("p50", (613, 820, 1753, 2352, 3355, 3660, 4125))
chart_data.add_series("p99", (1104, 2251, 3341, 3912, 4913, 5215, 6115))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS, Inches(0.82), Inches(3.65), Inches(7.15), Inches(2.62), chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.value_axis.maximum_scale = 6600
chart.value_axis.minimum_scale = 0
chart.value_axis.major_unit = 2000
chart.value_axis.has_major_gridlines = True
chart.category_axis.has_title = True
chart.category_axis.axis_title.text_frame.text = "hop count"
chart.value_axis.has_title = True
chart.value_axis.axis_title.text_frame.text = "latency (ms)"
chart.series[0].format.line.color.rgb = rgb(BLUE)
chart.series[1].format.line.color.rgb = rgb(ROSE)
chart.series[1].format.line.dash_style = 4
label_box(slide, "Small-node footprint",
          "62 MiB idle node\n233 MiB under crossing\n≈384 MiB node + database + cache",
          8.45, 3.68, 3.95, 1.62, "EEF1F3", MIST, INK, 17)
label_box(slide, "Wire size",
          "155 B check-in\n220 B forum post\n243 B Bangla broadcast",
          8.45, 5.3, 3.95, 1.38, "EEF3F6", BLUE, BLUE, 14.5)
add_footer(slide, 14)


# 15. Takeaway
slide = add_slide("What this project establishes", "Takeaway", TEAL)
add_text(slide,
         "A forum can continue inside surviving domestic IP components without trusting one platform or one administrator.",
         0.95, 1.58, 11.45, 1.0, 25, INK, True, align=PP_ALIGN.CENTER)
points = [
    ("Local continuity", "Each reachable island can publish, read and moderate locally."),
    ("Verifiable exchange", "Signed bytes, receipts and audit copies make tampering visible."),
    ("Conditional crossing", "A trusted bridge reconnects islands only while it can reach both."),
]
for i, (head, body) in enumerate(points):
    x = 0.72 + i * 4.18
    color = [TEAL, BLUE, AMBER][i]
    label_box(slide, head, body, x, 3.05, 3.85, 1.65, WHITE, color, color, 17)
add_box(slide,
        "Next evidence needed: a real multi-ISP field trial with separate clocks, independent auditors and measured censor interference.",
        1.25, 5.45, 10.85, 0.88, "F3E8EB", ROSE, 18, True, INK)
add_text(slide, "No mesh claim. No absolute anonymity claim. No delivery claim when every path is gone.",
         1.3, 6.56, 10.75, 0.3, 14, MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 15)


prs.core_properties.title = "Islands of Reach"
prs.core_properties.subject = "Censorship-resistant federated forum for national internet blackouts"
prs.core_properties.author = "Islands of Reach project"
prs.core_properties.keywords = "federation, censorship resistance, audit log, internet blackout"
prs.save(OUT)
print(OUT)
